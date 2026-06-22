"""Text extraction from uploaded files (PDF, PPTX, XLSX, DOCX).

This module turns a file uploaded through Streamlit into plain text that
can be translated. It supports four formats and enforces a maximum file
size. Only text is extracted: images, charts and embedded media are
ignored, so an image-heavy PDF may yield little or no text.

DESIGN NOTES
============
- Extraction functions take a file-like object (what Streamlit's
  st.file_uploader returns) and return a string. They contain no
  Streamlit UI calls, so they stay testable and reusable.
- The single public entry point, extract_text(), dispatches on the
  file's MIME type and applies the size limit. It raises
  FileTooLargeError or UnsupportedFileError for the two expected error
  conditions, which the UI layer is responsible for catching and
  presenting to the user.

WHAT IS AND IS NOT EXTRACTED
============================
- DOCX: paragraphs and tables, read in their real document order so a
  table's text stays near the paragraphs around it. Text boxes (text
  stored in drawing shapes) are not extracted.
- PPTX: text frames (titles, bullets, text boxes) and tables, slide by
  slide.
- XLSX: every cell of every sheet.
- PDF: the text layer of every page. Scanned, image-only PDFs have no
  text layer and yield little or nothing (no OCR is performed).

ABOUT EXTRACTED TEXT QUALITY
============================
Text pulled from PDFs and slide decks is often dense or out of order:
lines can run together, reading order may be wrong, headers and footers
may be interleaved with body text. This is inherent to those formats.
The translation pipeline's refinement/upgrade step is where layout is
cleaned up; extraction only aims to recover the words.
"""

from io import BytesIO

import fitz  # PyMuPDF
from docx import Document
from docx.document import Document as _DocxDocument
from docx.table import Table
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook
from pptx import Presentation


# Maximum accepted upload size. Kept deliberately modest: the extracted
# text is sent to a language model, so very large files mean high cost
# and long latency. Adjust here if the policy changes.
MAX_FILE_SIZE_MB = 20

# MIME types produced by st.file_uploader for each supported extension.
_MIME_PDF = "application/pdf"
_MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_MIME_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_MIME_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


class FileTooLargeError(Exception):
    """Raised when an uploaded file exceeds MAX_FILE_SIZE_MB."""


class UnsupportedFileError(Exception):
    """Raised when an uploaded file's type is not one we can read."""


# ---------------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------------

def read_pdf(file) -> str:
    """Extract text from a PDF file.

    Args:
        file: A file-like object (e.g. the object returned by
            st.file_uploader). Its bytes are read into memory.

    Returns:
        The concatenated text of every page. Pages that contain only
        images contribute nothing.
    """
    text_parts = []
    bytes_stream = BytesIO(file.read())
    with fitz.open(stream=bytes_stream, filetype="pdf") as document:
        for page in document:
            text_parts.append(page.get_text())
    return "".join(text_parts)


def _table_to_text(table: Table) -> str:
    """Flatten a Word table into text, one line per row.

    Cells within a row are separated by tabs; rows are separated by
    newlines. This keeps the tabular content readable and groups each
    row's cells together.

    Args:
        table: A python-docx Table object.

    Returns:
        The table's text. Empty cells contribute an empty string.
    """
    row_lines = []
    for row in table.rows:
        cells_text = [cell.text.strip() for cell in row.cells]
        row_lines.append("\t".join(cells_text))
    return "\n".join(row_lines)


def _iter_docx_block_text(document: _DocxDocument):
    """Yield the text of a Word document's blocks in reading order.

    python-docx exposes paragraphs and tables as two separate
    collections, which loses their relative order. To preserve reading
    order, this helper walks the document body's underlying XML and
    yields paragraph and table text as each block appears.

    Args:
        document: A python-docx Document object.

    Yields:
        Text strings, one per non-empty paragraph and one per table, in
        the order they appear in the document body.
    """
    body = document.element.body
    for child in body.iterchildren():
        if child.tag.endswith("}p"):
            # Wrap the raw XML paragraph element back into a Paragraph so
            # we can use its convenient .text property.
            paragraph = Paragraph(child, document)
            if paragraph.text.strip():
                yield paragraph.text
        elif child.tag.endswith("}tbl"):
            table = Table(child, document)
            table_text = _table_to_text(table)
            if table_text.strip():
                yield table_text


def read_docx(file) -> str:
    """Extract text from a Word document, paragraphs and tables included.

    Paragraphs and tables are read in their real order in the document,
    so a table's content stays close to the surrounding paragraphs.
    Text boxes (text held in drawing shapes) are not extracted.

    Args:
        file: A file-like object for the .docx file.

    Returns:
        The document's text, with blocks separated by newlines.
    """
    document = Document(file)
    return "\n".join(_iter_docx_block_text(document))


def read_pptx(file) -> str:
    """Extract text from a PowerPoint presentation, tables included.

    For each slide, collects the text of every shape that has a text
    frame (titles, bullet points, text boxes) and the text of every
    table cell.

    Args:
        file: A file-like object for the .pptx file.

    Returns:
        The collected text. Shapes are separated by spaces; table rows
        by newlines and cells within a row by tabs.
    """
    text_parts = []
    presentation = Presentation(file)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells_text = [cell.text.strip() for cell in row.cells]
                    text_parts.append("\t".join(cells_text))
            elif shape.has_text_frame:
                text_parts.append(shape.text)
    return " ".join(text_parts)


def read_excel(file) -> str:
    """Extract cell values from an Excel workbook.

    Reads every sheet, every row, every cell, converting each value to a
    string. Empty cells become the string "None"; this is acceptable for
    translation input but worth knowing when inspecting output.

    Args:
        file: A file-like object for the .xlsx file.

    Returns:
        All cell values joined by spaces, in row-major order across all
        sheets.
    """
    text_parts = []
    workbook = load_workbook(filename=file)
    for sheet in workbook:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                text_parts.append(str(cell))
    return " ".join(text_parts)


# ---------------------------------------------------------------------------
# Public dispatch
# ---------------------------------------------------------------------------

# Maps each supported MIME type to its extractor function.
_EXTRACTORS = {
    _MIME_PDF: read_pdf,
    _MIME_PPTX: read_pptx,
    _MIME_XLSX: read_excel,
    _MIME_DOCX: read_docx,
}


def extract_text(uploaded_file) -> str:
    """Extract plain text from a Streamlit-uploaded file.

    This is the single entry point the UI should call. It checks the
    file size, selects the right extractor based on the file's MIME
    type, and returns the extracted text.

    Args:
        uploaded_file: The object returned by st.file_uploader. Must
            expose `.size` (bytes), `.type` (MIME string) and the usual
            file-like read interface.

    Returns:
        The extracted text. May be empty if the file contained no
        extractable text (e.g. a scanned, image-only PDF).

    Raises:
        FileTooLargeError: If the file exceeds MAX_FILE_SIZE_MB. The
            message is suitable for display to the user.
        UnsupportedFileError: If the file's MIME type is not one of the
            four supported formats. The message is suitable for display.
    """
    size_limit_bytes = MAX_FILE_SIZE_MB * 1024 * 1024
    if uploaded_file.size > size_limit_bytes:
        raise FileTooLargeError(
            f"File too large. The maximum size is {MAX_FILE_SIZE_MB} MB."
        )

    extractor = _EXTRACTORS.get(uploaded_file.type)
    if extractor is None:
        raise UnsupportedFileError(
            "Unsupported file type. Please upload a PDF, PPTX, XLSX or DOCX file."
        )

    return extractor(uploaded_file)
